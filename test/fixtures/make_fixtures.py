#!/usr/bin/env python3
"""Regenerate the RAG real-file test fixtures used by test/UI_CHECKS.md.

These exercise the v0.67e structured extractors (items 6/7) and the embed
pipeline (item 8) on REAL bytes. All content is synthetic (no real data).

Usage:  pip install reportlab python-docx python-pptx openpyxl
        python3 make_fixtures.py
Outputs (this folder): test-report.pdf, test-policy.docx, test-slides.pptx, test-data.xlsx
"""
import os
HERE = os.path.dirname(os.path.abspath(__file__))

def pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(os.path.join(HERE, 'test-report.pdf'), pagesize=A4)
    story = []
    for h, b in [
        ('1. Executive Summary', 'Q4 revenue reached SGD 4.2 million, up 18% year on year. Net margin improved to 22%.'),
        ('2. Data Retention Policy', 'Operational logs are retained for 90 days. Personal data is purged after 24 months per PDPA.'),
        ('3. Security Posture', 'All endpoints run EDR. Quarterly penetration tests are mandated for internet-facing systems.'),
        ('3.1 Incident Response', 'Sev-1 incidents require notification to the CISO within 30 minutes.'),
    ]:
        story.append(Paragraph(h, styles['Heading2']))
        story.append(Paragraph(b, styles['BodyText']))
        story.append(Spacer(1, 12))
    doc.build(story)

def docx_():
    import docx
    d = docx.Document()
    d.add_heading('Employee Handbook', 0)
    d.add_heading('Leave Policy', level=1)
    d.add_paragraph('Annual leave is 21 days. Unused leave carries over up to 7 days.')
    d.add_heading('Remote Work', level=1)
    d.add_paragraph('Staff may work remotely up to 3 days per week with manager approval.')
    d.save(os.path.join(HERE, 'test-policy.docx'))

def pptx_():
    from pptx import Presentation
    p = Presentation()
    for title, body, note in [
        ('Project Overlord', 'Kickoff and scope', 'Internal only - do not distribute'),
        ('Timeline', 'Phase 1 Q1, Phase 2 Q3', 'Phase 2 depends on vendor signoff'),
        ('Budget', 'SGD 1.2M allocated', 'Contingency 10%'),
    ]:
        s = p.slides.add_slide(p.slide_layouts[1])
        s.shapes.title.text = title
        s.placeholders[1].text = body
        s.notes_slide.notes_text_frame.text = note
    p.save(os.path.join(HERE, 'test-slides.pptx'))

def xlsx_():
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active; ws.title = 'Q4 Sales'
    ws.append(['Region', 'Units', 'Revenue'])
    for r in [['APAC', 1200, 480000], ['EMEA', 800, 320000], ['Americas', 1500, 600000]]:
        ws.append(r)
    ws2 = wb.create_sheet('Headcount')
    ws2.append(['Dept', 'Staff']); ws2.append(['Engineering', 42]); ws2.append(['Sales', 18])
    wb.save(os.path.join(HERE, 'test-data.xlsx'))

if __name__ == '__main__':
    pdf(); docx_(); pptx_(); xlsx_()
    for f in ['test-report.pdf', 'test-policy.docx', 'test-slides.pptx', 'test-data.xlsx']:
        print(f, os.path.getsize(os.path.join(HERE, f)), 'bytes')
